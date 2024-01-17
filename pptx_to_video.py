from pptx import Presentation
from pptx.util import Inches
import subprocess
import fitz  # PyMuPDF
import cv2
import os
from moviepy.editor import ImageSequenceClip
from moviepy.video.compositing.concatenate import concatenate_videoclips

#importing API key etc...
import openai
import os
from openai import OpenAI
from dotenv import load_dotenv

load_dotenv()
openai_api = os.getenv('OPENAI_API')

if openai_api is None:
    raise ValueError("The OPEN_AI environment variable is not set.")
#don't change this part of the code. Should remain the same

#imports for slide creation
from pptx import Presentation
import re

#generate chatgpt slides content LLM content generation (prompt argument is the topic of the presentation)
def create_content(prompt):
    print("Converting to educational content....")
    client = openai.OpenAI(api_key=openai_api)
    story = client.chat.completions.create(
        messages = [
        {"role": "system", "content": "Create content for a 20-slide presentation on the input topic. For each slide, provide a title and up to six bullet points that effectively cover different aspects of the topic. Ensure that the content is informative, accurate, and well-structured to facilitate a comprehensive understanding of the subject. Lable the title of each slide as Title: followed by the name of the title of the slide"},
        {"role": "user", "content": prompt}],
        model="gpt-3.5-turbo",
    )
    # write response to the file
    nextpass = story.choices[0].message.content  
    file_name = "chatgptresponse.txt"
    with open(file_name, 'w') as file:
        # Write the string to the file
        file.write(nextpass)



#create presentation based on chatgpt text (stored in a text file)
def create_presentation_from_text(file_path):
    # Create a presentation object
    prs = Presentation()

    # Read the text file
    with open(file_path, 'r') as file:
        content = file.read()

    # Split the content into slides
    slides = re.split(r'Title:', content)[1:]  # Skip the first empty split

    for slide in slides:
        # Find the title and bullet points in the slide
        title_match = re.match(r'(.+)', slide)
        bullet_points_match = re.findall(r'- (.+)', slide)

        if title_match:
            # Create a new slide
            slide_layout = prs.slide_layouts[1]  # Using layout 1 for title and content
            new_slide = prs.slides.add_slide(slide_layout)
            title, subtitle = new_slide.shapes.title, new_slide.placeholders[1]

            # Set title
            title.text = title_match.group(1).strip()

            # Add bullet points
            for point in bullet_points_match:
                p = subtitle.text_frame.add_paragraph()
                p.text = point
                p.level = 0

    # Save the presentation
    prs.save('presentation.pptx')

    
def pdf_to_images(pdf_path, output_folder):
    pdf_document = fitz.open(pdf_path)
    
    for page_num in range(pdf_document.page_count):
        page = pdf_document[page_num]
        image = page.get_pixmap()
        
        image_path = f"{output_folder}/page_{page_num + 1}.png"
        image.save(image_path)

    pdf_document.close()

    
def images_to_video(image_folder, output_video_path, frame_rate=24):
    images = [img for img in os.listdir(image_folder) if img.endswith(".png")]

    def get_page_number(image_name):
        # Extract the numerical part of the filename
        return int(''.join(filter(str.isdigit, image_name)))

    images.sort(key=get_page_number)

    frame = cv2.imread(os.path.join(image_folder, images[0]))
    height, width, layers = frame.shape

    video = cv2.VideoWriter(output_video_path, cv2.VideoWriter_fourcc(*'mp4v'), frame_rate, (width, height))

    for image in images:
        img_path = os.path.join(image_folder, image)
        frame = cv2.imread(img_path)
        video.write(frame)

    cv2.destroyAllWindows()
    video.release()
    
def create_slideshow(images_folder, output_video_path, frame_rate=24, duration_per_image=5, transition_duration=1):
    images = [img for img in os.listdir(images_folder) if img.endswith(".png")]
    def get_page_number(image_name):
        # Extract the numerical part of the filename
        return int(''.join(filter(str.isdigit, image_name)))

    images.sort(key=get_page_number)

    clips = []
    for i in range(len(images)-1):
        current_img = os.path.join(images_folder, images[i])
        next_img = os.path.join(images_folder, images[i+1])

        transition_clip = ImageSequenceClip([current_img, next_img], fps=1/duration_per_image)
        # transition_clip = transition_clip.set_duration(duration_per_image)
        # transition_clip = transition_clip.fadein(transition_duration).fadeout(transition_duration)
        transition_clip = transition_clip.crossfadein(transition_duration)
        clips.append(transition_clip)

    final_clip = concatenate_videoclips(clips)

    final_clip.write_videofile(output_video_path, codec='libx264', fps=frame_rate)

if __name__ == "__main__":
    create_content("cheese") #generate chatgpt script for presentation - parameter is topic
    #change prompt as you feel is right in order to get a script that will flow well 
    file_path = "chatgptresponse.txt"
    create_presentation_from_text(file_path)# create slides from script file 
   
    
    # Example command: print the current working directory
    command = "ppt2pdf file presentation.pptx"  # On Windows, you can use "cd" instead

    # Execute the command and capture the result
    result = subprocess.run(command, shell=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)

    pdf_path = "test.pdf"
    images_folder = "slide_images"


    pdf_to_images(pdf_path, images_folder)
    output_video_path = "output_video.mp4"
    frame_rate = 24  # Adjust the frame rate as needed
    # images_to_video(image_folder, output_video_path, frame_rate)
    
    duration_per_image = 5  # Time each image is shown in seconds
    transition_duration = 1  # Duration of the transition between images in seconds
    create_slideshow(images_folder, output_video_path, frame_rate, duration_per_image, transition_duration)
    
    