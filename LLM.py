

import openai
import os
from openai import OpenAI
from dotenv import load_dotenv
from pptx import Presentation
import re
import os
import comtypes.client

load_dotenv()
openai_api = os.getenv('OPENAI_API')

if openai_api is None:
    raise ValueError("The OPEN_AI environment variable is not set.")

#generate chatgpt slides content LLM content generation (prompt argument is the topic of the presentation)
def create_content(prompt):
    print("Converting to educational content....")
    client = openai.OpenAI(api_key=openai_api)
    story = client.chat.completions.create(
        messages = [
        {"role": "system", "content": "Create content for a 20-slide presentation on the input topic. For each slide, provide a title and up to six bullet points that effectively cover different aspects of the topic. Ensure that the content is informative, accurate, and well-structured to facilitate a comprehensive understanding of the subject. Lable the title of each slide as Title: followed by the name of the title of the slide"},
        {"role": "user", "content": prompt}],
        model="gpt-3.5-turbo",
    )
    # write response to the file
    nextpass = story.choices[0].message.content  
    file_name = "chatgptresponse.txt"
    with open(file_name, 'w') as file:
        # Write the string to the file
        file.write(nextpass)



def create_presentation_from_text(file_path,pptx_file):
    # Create a presentation object
    prs = Presentation()

    # Read the text file
    with open(file_path, 'r') as file:
        content = file.read()

    # Split the content into slides
    slides = re.split(r'Title:', content)[1:]  # Skip the first empty split

    for slide in slides:
        # Find the title and bullet points in the slide
        title_match = re.match(r'(.+)', slide)
        bullet_points_match = re.findall(r'- (.+)', slide)

        if title_match:
            # Create a new slide
            slide_layout = prs.slide_layouts[1]  # Using layout 1 for title and content
            new_slide = prs.slides.add_slide(slide_layout)
            title, subtitle = new_slide.shapes.title, new_slide.placeholders[1]

            # Set title
            title.text = title_match.group(1).strip()

            # Add bullet points
            for point in bullet_points_match:
                p = subtitle.text_frame.add_paragraph()
                p.text = point
                p.level = 0

    # Save the presentation
    prs.save(pptx_file)


def pptx_to_pdf(pptx_filename, pdf_filename):
    # Path to the PowerPoint application
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")

    # PowerPoint must be visible to execute the conversion
    powerpoint.Visible = 1

    # Open the PowerPoint file
    presentation = powerpoint.Presentations.Open(pptx_filename)

    # Convert to PDF
    presentation.SaveAs(pdf_filename, FileFormat=32) # 32 for PDF format

    # Close the presentation and PowerPoint
    presentation.Close()
    powerpoint.Quit()


create_content("cheese")
file_path = "chatgptresponse.txt"
pptx_file = 'presentation.pptx'  # Path to your PowerPoint file
create_presentation_from_text(file_path,pptx_file)
pdf_file = 'presentation.pdf'    # Desired path for the PDF file
# Convert the file
pptx_to_pdf(pptx_file, pdf_file)


